#!/usr/bin/env python
"""Convert docs/presentation/index.html (reveal.js deck) into a native, editable
PowerPoint file (.pptx) suitable for import & editing in Google Slides.

The output uses text boxes, rounded rectangles, pictures and speaker notes so
every element remains editable rather than being flattened into screenshots.

Usage:
    python build_pptx.py [../FOSS4G-tidal-oss.pptx]

Requires: python-pptx, beautifulsoup4, lxml, pillow
"""

import os
import re
import sys
import html as html_mod

from bs4 import BeautifulSoup

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = os.path.dirname(os.path.abspath(__file__))
PRESENTATION = os.path.join(BASE, "index.html")
ASSETS = os.path.join(BASE, "assets")
PNG_CACHE = os.path.join(ASSETS, "png")
OUT_DEFAULT = os.path.join(BASE, "FOSS4G-tidal-oss.pptx")

# ---------------------------------------------------------------------------
# Theme (from the reveal deck's CSS :root block)
# ---------------------------------------------------------------------------
DEEP = RGBColor(0x0B, 0x3C, 0x5D)
OCEAN = RGBColor(0x0E, 0x74, 0x90)
AQUA = RGBColor(0x1A, 0xA7, 0xC4)
CORAL = RGBColor(0xE7, 0x6F, 0x51)
AMBER = RGBColor(0xF4, 0xA2, 0x61)
INK = RGBColor(0x12, 0x23, 0x2E)
PAPER = RGBColor(0xFB, 0xFD, 0xFF)
MUTED = RGBColor(0x5A, 0x70, 0x7E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GHOST = RGBColor(0xE6, 0xEE, 0xF2)
CARD_BORDER = RGBColor(0xD9, 0xE6, 0xEC)
CARD_BG = RGBColor(0xF4, 0xFA, 0xFC)
CARD_CORAL_BG = RGBColor(0xFD, 0xF5, 0xF1)
CARD_AMBER_BG = RGBColor(0xFD, 0xF9, 0xF0)
FORMULA_BG = RGBColor(0xED, 0xF6, 0xF9)
REFA_CLR = RGBColor(0x93, 0xA7, 0xB3)
AMBER_TEXT = RGBColor(0x3A, 0x2C, 0x00)

SANS = "Segoe UI"
SERIF = "Georgia"

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.55
TOP = 0.42
CONTENT_W = SLIDE_W - 2 * MARGIN
BODY_TOP = 1.30

# ---------------------------------------------------------------------------
# Slide size / presentation bootstrap
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------------------
# Low-level text helpers
# ---------------------------------------------------------------------------
def _no_autofit(tf):
    """Set autofit off for a text frame (keeps text exactly where placed)."""
    bodyPr = tf._txBody.find(qn('a:bodyPr'))  # noqa: SLF001
    for tag in ('a:normAutofit', 'a:spAutoFit'):
        el = bodyPr.find(qn(tag))
        if el is not None:
            bodyPr.remove(el)
    bodyPr.set('wrap', 'square')


def add_text(slide, x, y, w, h, text, size, color=INK, bold=False, italic=False,
             align=PP_ALIGN.LEFT, font=SANS, line_spacing=1.0, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    _no_autofit(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.color.rgb = color
    f.bold = bold
    f.italic = italic
    return tb


def add_runs(slide, x, y, w, h, runs, size, align=PP_ALIGN.LEFT, color=INK,
             line_spacing=1.0, anchor=MSO_ANCHOR.TOP):
    """runs = list of (text, bold) tuples rendered on a single line."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    _no_autofit(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for text, bold in runs:
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = SANS
        f.size = Pt(size)
        f.color.rgb = color
        f.bold = bold
    return tb


# ---------------------------------------------------------------------------
# Shape helpers
# ---------------------------------------------------------------------------
def add_rect(slide, x, y, w, h, fill, line=None, shadow=False,
             shape=MSO_SHAPE.RECTANGLE, radius=None):
    shp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.4)
    shp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    return shp


def add_rule(slide, x, y, w, colors):
    """Gradient-like accent bar: render as stacked colour segments."""
    n = len(colors)
    seg = w / n
    for i, c in enumerate(colors):
        add_rect(slide, x + i * seg, y, seg, 0.09, c)


def _fit_left(shape):
    pass


def autosize_card_height(slide, x, y, w, title, body_lines, size=11.0,
                         title_size=13.0, pad=0.16, gap=0.05):
    """Roughly estimate the height needed for a card body so text fits."""
    if not body_lines:
        body_lines = [""]
    lines_per_row = compute_lines(w - 0.34, size)
    rows = 0
    for bl in body_lines:
        rows += est_lines(bl, lines_per_row)
    h = pad * 2 + 0.24 + rows * (size / 72.0 + gap) + 0.02
    return max(0.55, h)


def compute_lines(width_in, font_pt):
    """Average chars that fit in a box of given width at given font size."""
    chars = max(4, int(width_in * 72 / (font_pt * 0.50)))
    return chars


def est_lines(text, chars_per_row):
    if not text:
        return 1
    manual = text.count("\n")
    lines = 0
    for seg in text.split("\n"):
        n = (len(seg) // chars_per_row) + (1 if len(seg) % chars_per_row else 0)
        lines += max(1, n)
    return lines + manual


def add_card(slide, x, y, w, title, body_lines, variant="", big_title=False):
    border = CARD_BORDER
    bg = CARD_BG
    accent = OCEAN
    title_size = 13.5 if big_title else 13.0
    size = 12.0 if big_title else 11.0
    if variant == "coral":
        bg, accent = CARD_CORAL_BG, CORAL
    elif variant == "amber":
        bg, accent = CARD_AMBER_BG, AMBER

    h = autosize_card_height(slide, x, y, w, title, body_lines,
                             size=11.0 if not big_title else 12.0)
    # body
    add_rect(slide, x, y, w, h, bg, line=border, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.06)
    # left accent bar
    add_rect(slide, x, y, 0.075, h, accent)
    # title
    add_text(slide, x + 0.22, y + 0.14, w - 0.36, 0.26, title,
             title_size, DEEP, bold=True)
    # body lines
    yy = y + 0.40
    for bl in body_lines:
        if bl.strip() == "":
            continue
        add_text(slide, x + 0.22, yy, w - 0.36, 0.14, bl, size, INK,
                 line_spacing=1.02)
        yylines = est_lines(bl, compute_lines(w - 0.36, size))
        yy += yylines * (size / 72.0 + 0.05)


def add_chip(slide, x, y, text, variant="", size=10.0):
    w = 0.075 * len(text) + 0.30
    bg = DEEP
    fg = WHITE
    if variant == "ocean":
        bg = OCEAN
    elif variant == "coral":
        bg = CORAL
    elif variant == "amber":
        bg = AMBER
        fg = AMBER_TEXT
    elif variant == "ghost":
        bg = GHOST
        fg = DEEP
    add_rect(slide, x, y, w, 0.30, bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    add_text(slide, x, y + 0.045, w, 0.22, text, size, fg, bold=True,
             align=PP_ALIGN.CENTER)
    return w


def add_stat(slide, x, y, w, h, strong, label):
    add_rect(slide, x, y, w, h, DEEP, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    add_text(slide, x, y + 0.10, w, 0.34, strong, 17.0, RGBColor(0x8D, 0xE7, 0xF2),
             bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + 0.50, w, 0.34, label, 9.5, WHITE, align=PP_ALIGN.CENTER)


def add_stat_band(slide, x, y, w, stats, cols=None):
    n = len(stats)
    cols = cols or n
    gap = 0.22
    cw = (w - gap * (cols - 1)) / cols
    h = 1.03
    for i, st in enumerate(stats):
        row = i // cols
        col = i % cols
        add_stat(slide, x + col * (cw + gap), y + row * (h + 0.12), cw, h,
                 st[0], st[1])
    return h + (rows_of(cols, n) - 1) * (h + 0.12)


def rows_of(cols, n):
    return max(1, -(-n // cols))


def add_formula(slide, x, y, w, h, text, size=14.0):
    add_rect(slide, x, y, w, h, FORMULA_BG, line=AQUA, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.10)
    add_text(slide, x, y + h * 0.14, w, h * 0.72, text, size, DEEP, italic=True,
             align=PP_ALIGN.CENTER, font=SERIF)


def add_blockquote(slide, x, y, w, text, cite="", size=12.0):
    # determine height from text
    lines = est_lines(text, compute_lines(w - 0.5, size))
    h = 0.34 + lines * (size / 72.0 + 0.05) + (0.34 if cite else 0)
    add_rect(slide, x, y, w, h, CARD_CORAL_BG, line=CORAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.05)
    add_text(slide, x + 0.22, y + 0.12, w - 0.44, 0.2, "\u201c%s\u201d" % text, size,
             DEEP, italic=True, line_spacing=1.05)
    if cite:
        add_text(slide, x + 0.22, y + h - 0.30, w - 0.44, 0.24, cite, size * 0.85,
                 MUTED)
    return h


# ---------------------------------------------------------------------------
# Image helpers
# ---------------------------------------------------------------------------
def resolve_image(src):
    """Map an <img src> to a real PNG path on disk."""
    base = os.path.basename(src)
    stem, ext = os.path.splitext(base)
    if ext.lower() == ".svg":
        cand = os.path.join(PNG_CACHE, stem + ".png")
        if os.path.exists(cand):
            return cand
        return os.path.join(PNG_CACHE, base.replace(".svg", ".png"))
    # already a png: prefer cache copy if present else asset dir
    cand = os.path.join(PNG_CACHE, base)
    if os.path.exists(cand):
        return cand
    return os.path.join(ASSETS, base)


def add_picture_fit(slide, img_path, x, y, max_w, max_h, align=PP_ALIGN.CENTER):
    from PIL import Image
    with Image.open(img_path) as im:
        pw, ph = im.size
    ar = pw / ph
    w = max_w
    h = w / ar
    if h > max_h:
        h = max_h
        w = h * ar
    ix = x + (max_w - w) / 2 if align == PP_ALIGN.CENTER else x
    pic = slide.shapes.add_picture(img_path, Inches(ix), Inches(y), Inches(w), Inches(h))
    return pic, w, h


def figure_block(section, slide, x, y, w, img, caption=""):
    """Place a picture fit to (w, ~3.6in) with optional caption below."""
    pic, pw, ph = add_picture_fit(slide, resolve_image(img), x, y, w, 3.55)
    if caption:
        cy = y + ph + 0.10
        add_text(slide, x, cy, w, 0.30, caption, 9.5, MUTED, align=PP_ALIGN.CENTER)
        return ph + 0.40
    return ph + 0.12


# ---------------------------------------------------------------------------
# HTML traversal helpers
# ---------------------------------------------------------------------------
def node_text(node):
    return re.sub(r"\s+", " ", node.get_text(" ", strip=True)).strip()


def node_text_preserve(node):
    """Text that preserves <br> as newlines and trims internal spaces."""
    parts = []
    for child in node.descendants:
        if getattr(child, "name", None) == "br":
            parts.append("\n")
        elif getattr(child, "name", None) is None:
            parts.append(str(child))
    txt = "".join(parts)
    return txt


def clean_ws(t):
    t = t.replace("\u00a0", " ")
    lines = t.split("\n")
    return "\n".join(re.sub(r"[\x00-\x1f]+", "", re.sub(r"[ \t]+", " ", l)).strip()
                     for l in lines).strip()


def parse_color_runs(node):
    """Return list of (text, bold) for a <p>/<div> preserving <b>/<strong>."""
    runs = []
    for child in node.descendants:
        if getattr(child, "name", None) in ("br", None):
            txt = str(child).strip() if getattr(child, "name") is None else "\n"
            if child is None:
                continue
        else:
            continue
    return runs


def is_chips_row(node):
    kids = [c for c in getattr(node, "children", []) if getattr(c, "name", None)]
    if not kids:
        return False
    return all(c.get("class") and "chip" in c.get("class") for c in kids)


def chip_variant(el):
    cls = el.get("class") or []
    for v in ("ocean", "coral", "amber", "ghost"):
        if v in cls:
            return v
    return ""


def card_variant(el):
    cls = el.get("class") or []
    for v in ("coral", "amber"):
        if v in cls:
            return v
    return ""


def collect_card(el):
    """Extract (title, body_lines, variant) from a .card element."""
    variant = card_variant(el)
    title = ""
    body_lines = []
    for child in el.find_all(["h4", "p"], recursive=True):
        if child.parent != el and child.name != "p":
            # only take direct h4
            if child.name == "h4" and child.parent != el:
                continue
        if child.name == "h4":
            title = node_text(child)
        elif child.name == "p":
            txt = node_text(child)
            if txt:
                body_lines.append(txt)
    if not body_lines and title:
        # maybe h4 + p nested (p containing .chip spans)
        pass
    return title, body_lines, variant


# ---------------------------------------------------------------------------
# Core per-section layout
# ---------------------------------------------------------------------------
def est_card_h(w, title, body_lines):
    size = 11.0
    lines = 0
    for bl in body_lines:
        lines += est_lines(bl, compute_lines(w - 0.36, size))
    return 0.40 + lines * (size / 72.0 + 0.05) + 0.16


def layout(slide, sect):
    return _layout(slide, sect)


def _layout(slide, sect):
    top = TOP
    used_header = False
    # Handle header
    title_el = sect.find(["h1", "h2"])
    kick_el = sect.find("h3")
    kick_el = kick_el if kick_el else sect.find(class_="kick")
    subtitle_el = sect.find(class_="subtitle")
    if not subtitle_el:
        subtitle_el = sect.find("p", class_="kick")

    kicker = node_text(kick_el) if kick_el else ""
    title = node_text(title_el) if title_el else ""
    subtitle = node_text(subtitle_el) if subtitle_el else ""

    # --- render header ---
    header_center = False
    if sect.get("class") and "title-slide" in sect.get("class"):
        header_center = True
    y = top
    if kicker:
        if header_center:
            add_text(slide, MARGIN, y, CONTENT_W, 0.30, kicker.upper(), 11.5, OCEAN,
                     bold=True, align=PP_ALIGN.CENTER)
        else:
            add_text(slide, MARGIN, y, CONTENT_W, 0.28, kicker.upper(), 11.5, OCEAN,
                     bold=True)
        y += 0.36
    if title:
        tsize = 27 if header_center else 24
        tl = est_lines(title, compute_lines(CONTENT_W, tsize))
        add_text(slide, MARGIN, y, CONTENT_W, 0.4, title, tsize, DEEP, bold=True,
                 align=PP_ALIGN.CENTER if header_center else PP_ALIGN.LEFT)
        y += 0.52
    # title-slide gradient rule
    if header_center:
        add_rule(slide, SLIDE_W / 2 - 0.6, y, 1.2, [OCEAN, AQUA, CORAL])
        y += 0.24
    if subtitle:
        sub_size = 13.5
        sll = est_lines(subtitle, compute_lines(CONTENT_W, sub_size))
        add_text(slide, MARGIN, y, CONTENT_W, 0.4, subtitle, sub_size, MUTED,
                 align=PP_ALIGN.CENTER if header_center else PP_ALIGN.LEFT,
                 line_spacing=1.05)
        y += 0.34

    body_y = max(y, BODY_TOP)

    # Walk child blocks in document order, rendering each row.
    footer_ref = None
    footnote = None

    for child in sect.children:
        if getattr(child, "name", None) is None:
            continue
        cls = child.get("class") or []
        # skip the header elements already rendered
        if child is title_el or child is kick_el or child is subtitle_el:
            continue
        if child.name == "aside":
            continue
        name = child.name

        # --- full-width-ish blocks handled inline ---
        if "ref" in cls:
            footer_ref = node_text(child)
            continue
        if "footnote" in cls or (name == "p" and "footnote" in cls):
            footnote = node_text(child)
            continue

        if name == "div" and "stat-band" in cls:
            stats = []
            for st in child.find_all("div", class_="stat", recursive=False):
                strong = st.find("strong")
                span = st.find("span")
                stats.append((node_text(strong) if strong else "",
                              node_text(span) if span else ""))
            add_stat_band(slide, MARGIN, body_y, CONTENT_W, stats)
            body_y += 1.15
            continue

        if name == "div" and "compare" in cls:
            h = render_compare(slide, child, MARGIN, body_y, CONTENT_W)
            body_y += h + 0.15
            continue

        if name == "blockquote":
            text = node_text(child)
            cite_el = child.find("cite")
            cite_text = node_text(cite_el) if cite_el else ""
            h = add_blockquote(slide, MARGIN, body_y, CONTENT_W, text, cite_text)
            body_y += h + 0.12
            continue

        if name == "div" and "datapath" in cls:
            h = render_datapath(slide, child, MARGIN, body_y, CONTENT_W)
            body_y += h + 0.10
            continue

        if name in ("div", "p") and ("formula" in cls):
            ftxt = node_text(child)
            add_formula(slide, MARGIN, body_y, CONTENT_W, 0.62, ftxt)
            body_y += 0.74
            continue

        # grid columns / hotspot-detail / webgis / visual-pair / visual-grid
        if "grid" in cls or "hotspot-detail" in cls:
            h = render_grid(slide, child, MARGIN, body_y, CONTENT_W)
            body_y += h + 0.14
            continue

        # visual-panel with a picture
        if "visual-panel" in cls:
            h = render_visual_panel(slide, child, MARGIN, body_y, CONTENT_W)
            body_y += h + 0.12
            continue

        # bare svg (full-stack slide)
        if name == "svg":
            h = render_svg(slide, child, MARGIN, body_y, CONTENT_W)
            body_y += h + 0.12
            continue

        # bare image figure (img directly under section)
        if name == "img":
            cap_el = sect and caption_for_img(sect, child)
            h = figure_block(sect, slide, MARGIN, body_y, CONTENT_W, child.get("src"),
                             cap_el)
            body_y += h + 0.12
            continue

    # footer: ref + footnote
    if footnote or footer_ref:
        footer_y = SLIDE_H - 0.42
        if footer_ref:
            add_text(slide, SLIDE_W - MARGIN - 6.2, footer_y, 6.2, 0.26, footer_ref,
                     8.5, MUTED, align=PP_ALIGN.RIGHT)
        if footnote:
            add_text(slide, MARGIN, footer_y, CONTENT_W - 6.0, 0.26, footnote, 8.5,
                     MUTED)


def caption_for_img(sect, img_el):
    # look for a following sibling div.fig-cap
    nxt = img_el.find_next_sibling()
    if nxt and nxt.get("class") and "fig-cap" in nxt.get("class"):
        return node_text(nxt)
    return ""


def render_datapath(slide, el, x, y, w):
    """Render a .datapath row: chips separated by '->' arrows."""
    tx = x
    for child in el.children:
        if getattr(child, "name", None) is None:
            continue
        if child.name in ("b", "span") and child.name == "b":
            pass
        if child.name == "span" and child.get("class") and "chip" in child.get("class"):
            v = chip_variant(child)
            cw = add_chip(slide, tx, y, node_text(child), v)
            tx += cw + 0.14
        elif child.name in ("b",):
            add_text(slide, tx, y + 0.02, 0.3, 0.26, "\u2192", 14, CORAL, bold=True,
                     align=PP_ALIGN.CENTER)
            tx += 0.20
    return 0.30


def render_svg(slide, svg_el, x, y, w):
    label = ""
    al = svg_el.get("aria-label") or ""
    img = None
    if "FOSS4G stack" in al or "End-to-end" in al:
        img = os.path.join(PNG_CACHE, "full_stack.png")
    elif "MSP tool" in al or "layered architecture" in al:
        img = os.path.join(PNG_CACHE, "msp_architecture.png")
    if not img or not os.path.exists(img):
        return 0.0
    pic, pw, ph = add_picture_fit(slide, img, x, y, w, 4.4)
    return ph


def render_visual_panel(slide, el, x, y, w, box_h=None):
    """A .visual-panel: bordered box holding an <img> / inline svg / datapath
    chip row / plain text row, plus an optional label. box_h controls the media
    height; if not given it is derived from the panel class."""
    cls = el.get("class") or []
    if box_h is None:
        if "mcda-panel" in cls:
            box_h = 1.55
        elif "roadmap-panel" in cls:
            box_h = 1.65
        elif "timeseries-panel" in cls:
            box_h = 1.35
        elif "diagram-panel" in cls:
            box_h = 3.0
        elif "turbine-visuals" in cls:
            box_h = 2.4
        else:
            box_h = 3.0
    img_el = el.find("img")
    svg_el = el.find("svg")
    datapath_el = el.find(class_="datapath")
    label_el = el.find(class_="visual-label")
    label = node_text(label_el) if label_el else ""
    # the bare text row (e.g. RESOURCE -> CONSTRAINTS -> SUITABILITY -> INVESTMENT)
    text_row = panel_text_row(el)

    # --- datapath panel: label on top, chips row beneath ---
    if datapath_el and not img_el and not svg_el:
        panel_h = 0.30 + (0.30 if label else 0.04) + 0.14
        add_rect(slide, x, y, w, panel_h, WHITE, line=CARD_BORDER,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
        if label:
            add_text(slide, x + 0.16, y + 0.10, w - 0.32, 0.28, label, 9.5, MUTED)
        render_datapath(slide, datapath_el, x + 0.16, y + 0.34, w - 0.32)
        return panel_h

    # --- text-only panel (label + one-row of labelled arrows/words) ---
    if not img_el and not svg_el and text_row:
        rows = est_lines(text_row, compute_lines(w - 0.5, 10.5))
        panel_h = 0.30 + (0.28 if label else 0.04) + rows * (10.5 / 72.0 + 0.04)
        add_rect(slide, x, y, w, panel_h, WHITE, line=CARD_BORDER,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
        if label:
            add_text(slide, x + 0.20, y + 0.10, w - 0.40, 0.28, label, 9.5, MUTED)
        add_text(slide, x + 0.20, y + 0.36, w - 0.40, 0.28, text_row, 10.5, OCEAN,
                 bold=True, align=PP_ALIGN.CENTER)
        return panel_h

    panel_h = 0.28 + box_h + (0.30 if label else 0.06)
    add_rect(slide, x, y, w, panel_h, WHITE, line=CARD_BORDER,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    if img_el:
        add_picture_fit(slide, resolve_image(img_el.get("src")),
                        x + 0.16, y + 0.14, w - 0.32, box_h)
    elif svg_el:
        lbl = svg_el.get("aria-label") or ""
        path = None
        if "MSP tool" in lbl or "layered" in lbl:
            path = os.path.join(PNG_CACHE, "msp_architecture.png")
        elif "FOSS4G stack" in lbl or "End-to-end" in lbl:
            path = os.path.join(PNG_CACHE, "full_stack.png")
        if path:
            add_picture_fit(slide, path, x + 0.16, y + 0.14, w - 0.32, box_h)
    if label:
        add_text(slide, x + 0.16, y + 0.18 + box_h, w - 0.32, 0.28, label, 9.5, MUTED,
                 align=PP_ALIGN.CENTER)
    return panel_h


def panel_text_row(el):
    """If a .visual-panel contains only a row of text (label + words/arrows)
    with no media, return that text, else ''."""
    if el.find("img") or el.find("svg") or el.find(class_="datapath"):
        return ""
    label_el = el.find(class_="visual-label")
    text = ""
    for child in el.children:
        if getattr(child, "name", None) is None:
            continue
        if child is label_el:
            continue
        txt = node_text(child)
        if txt:
            return txt
    return ""


CONTENT_CLASSES = ("card", "visual-panel", "datapath", "formula", "webgis-copy",
                   "webgis-figure", "hotspot-copy", "hotspot-visuals",
                   "visual-pair", "visual-grid", "compare", "stat-band",
                   "center", "mcda-panel", "roadmap-panel", "timeseries-panel",
                   "diagram-panel")


def render_grid(slide, el, x, y, w):
    cls = el.get("class") or []
    cols = 3 if "three-col" in cls else 2
    gap = 0.30

    if "hotspot-detail" in cls:
        return render_hotspot(slide, el, x, y, w)

    children = [c for c in el.children if getattr(c, "name", None) is not None]

    columns = []
    for c in children:
        if is_column_wrapper(c):
            columns.append([cc for cc in c.children
                            if getattr(cc, "name", None) is not None])
        else:
            # a standalone content element forms its own column
            columns.append([c])
    cols = max(cols, len(columns))

    cw = (w - gap * (cols - 1)) / cols
    max_h = 0
    for i, items in enumerate(columns):
        colx = x + i * (cw + gap)
        max_h = max(max_h, render_column(slide, items, colx, y, cw))
    return max_h


def is_column_wrapper(el):
    """A plain <div> that groups multiple content elements (grid column)."""
    if el.name != "div":
        return False
    if is_content_el(el):
        return False
    if el.get("class") is None and (el.get("style") is not None):
        return True
    # plain div (no content class) that itself has content descendants
    return el.find(["div", "img", "svg", "blockquote"]) is not None


def is_content_el(el):
    cls = el.get("class") or []
    if any(c in cls for c in CONTENT_CLASSES):
        return True
    if el.name in ("img", "svg", "blockquote"):
        return True
    return False


def distribute(children, cols):
    """Split a flat list of content elements into n columns round-robin."""
    columns = [[] for _ in range(cols)]
    for i, c in enumerate(children):
        columns[i % cols].append(c)
    return columns


def render_column(slide, items, x, y, w):
    yy = y
    for item in items:
        item_cls = item.get("class") or []
        name = item.name
        if name in ("h3", "h4"):
            continue
        if name == "div" and "card" in item_cls:
            title, body_lines, variant = collect_card(item)
            h = est_card_h(w, title, body_lines)
            add_card(slide, x, yy, w, title, body_lines, variant)
            yy += h + 0.12
        elif "visual-panel" in item_cls:
            h = render_visual_panel(slide, item, x, yy, w)
            yy += h + 0.12
        elif name == "div" and "visual-pair" in item_cls:
            yy += render_image_pair(slide, item, x, yy, w, max_h=2.1, labeled=True)
        elif name == "div" and "visual-grid" in item_cls:
            yy += render_image_pair(slide, item, x, yy, w, max_h=3.0, labeled=False)
        elif name == "div" and ("center" in item_cls or "webgis-copy" in item_cls):
            yy += render_nested_column(slide, item, x, yy, w)
        elif name == "div" and "webgis-figure" in item_cls:
            img_el = item.find("img")
            if img_el:
                cap_el = caption_for_img(item, img_el)
                ph = figure_block(item, slide, x, yy, w, img_el.get("src"), cap_el)
                yy += ph + 0.12
        elif name == "div" and "ref" in item_cls:
            continue  # handled as footer
        elif name in ("img", "svg"):
            ph = figure_block(None, slide, x, yy, w, item.get("src"), "")
            yy += ph + 0.12
    return yy - y


def render_nested_column(slide, el, x, y, w):
    """A column div containing .card elements (e.g. a plain wrapper or center)."""
    yy = y
    items = [c for c in el.children if getattr(c, "name", None) is not None]
    # if the wrapper has no card children but an img, treat as figure
    has_card = any("card" in (it.get("class") or []) for it in items)
    if not has_card:
        img_el = el.find("img") or None
        cap_el = caption_for_img(el, img_el) if img_el else ""
        if img_el:
            ph = figure_block(el, slide, x, yy, w, img_el.get("src"), cap_el)
            return ph + 0.12
        return 0
    for it in items:
        it_cls = it.get("class") or []
        if it.name == "div" and "card" in it_cls:
            title, body_lines, variant = collect_card(it)
            ch = est_card_h(w, title, body_lines)
            add_card(slide, x, yy, w, title, body_lines, variant)
            yy += ch + 0.12
    return yy - y


def render_image_pair(slide, el, x, y, w, max_h, labeled):
    """A .visual-pair / .visual-grid: side-by-side child panels (each holding an
    image + label), or bare images. Renders them as a 2-column sub-grid."""
    children = [c for c in el.children if getattr(c, "name", None) is not None]
    panels = [c for c in children if c.name == "div"
              and ("visual-panel" in (c.get("class") or []))]
    bare_imgs = [c for c in children if c.name == "img"]

    if panels:
        gap = 0.14
        sub_w = (w - gap) / 2
        ph = 0
        for k, p in enumerate(panels):
            sub_h = render_visual_panel(slide, p, x + k * (sub_w + gap), y,
                                        sub_w, box_h=max_h)
            ph = max(ph, sub_h)
        return ph + 0.12

    if not bare_imgs:
        return 0
    gap = 0.14
    sub_w = (w - gap) / 2
    ph = 0
    for k, im in enumerate(bare_imgs):
        p = add_picture_fit(slide, resolve_image(im.get("src")),
                            x + k * (sub_w + gap), y, sub_w, max_h)
        ph = max(ph, p[2])
    if labeled:
        labels = [node_text(l) for l in el.find_all(class_="visual-label")]
        if any(labels):
            add_text(slide, x, y + ph + 0.05, w, 0.28,
                     labels[0], 9.5, MUTED, align=PP_ALIGN.CENTER)
            ph += 0.32
    return ph + 0.12


def render_hotspot(slide, el, x, y, w):
    # .hotspot-copy (2-col grid of cards) + .hotspot-visuals (2 images)
    copy = el.find(class_="hotspot-copy")
    visuals = el.find(class_="hotspot-visuals")
    h = 0
    if copy:
        cards = copy.find_all("div", class_="card")
        left = cards[::2]
        right = cards[1::2]
        gapx = 0.30
        cw = (w - gapx) / 2
        maxh = 0
        for col_idx, card_list in enumerate((left, right)):
            colx = x + col_idx * (cw + gapx)
            yy = y
            for card in card_list:
                title, body_lines, variant = collect_card(card)
                ch = est_card_h(cw, title, body_lines)
                add_card(slide, colx, yy, cw, title, body_lines, variant)
                yy += ch + 0.12
            maxh = max(maxh, yy - y)
        h = maxh
    if visuals:
        imgs = visuals.find_all("img")
        gapx = 0.30
        cw = (w - gapx) / 2
        ph = 0
        yy = y + h + 0.14
        for k, im in enumerate(imgs):
            p = add_picture_fit(slide, resolve_image(im.get("src")),
                                x + k * (cw + gapx), yy, cw, 3.2)
            ph = max(ph, p[2])
        # fig-cap below
        cap_el = el.find(class_="fig-cap")
        if cap_el:
            add_text(slide, MARGIN, yy + ph + 0.08, w, 0.28, node_text(cap_el),
                     9.5, MUTED, align=PP_ALIGN.CENTER)
            ph += 0.34
        h += 0.14 + ph
    return h


def render_compare(slide, el, x, y, w):
    gap = 0.30
    cw = (w - gap) / 2
    h = 0
    parts = [c for c in el.children if getattr(c, "name", None) is not None]
    for i, part in enumerate(parts):
        cls = part.get("class") or []
        is_open = "open" in cls
        boxx = x + i * (cw + gap)
        # title
        h4 = part.find("h4")
        thead = node_text(h4) if h4 else ""
        ps = [node_text(p) for p in part.find_all("p")]
        # panel
        border = AQUA if is_open else CORAL
        bg = CARD_CORAL_BG if not is_open else CARD_BG
        titlec = OCEAN if is_open else CORAL
        lines = sum(est_lines(p, compute_lines(cw - 0.4, 11.0)) for p in ps)
        ph = 0.46 + lines * (11 / 72 + 0.05) + 0.16
        add_rect(slide, boxx, y, cw, ph, bg, line=border,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
        add_text(slide, boxx + 0.2, y + 0.13, cw - 0.4, 0.26, thead, 13.5, titlec,
                 bold=True)
        yy = y + 0.42
        for p in ps:
            add_text(slide, boxx + 0.2, yy, cw - 0.4, 0.2, p, 11, INK)
            yy += (11 / 72 + 0.05)
        h = max(h, ph)
    return h


# ---------------------------------------------------------------------------
# Speaker notes
# ---------------------------------------------------------------------------
def add_notes(slide, sect):
    aside = sect.find("aside", class_="notes")
    if not aside:
        return
    text = node_text_preserve(aside) or node_text(aside)
    text = clean_ws(text)
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main(out_path):
    with open(PRESENTATION, encoding="utf-8") as f:
        html = f.read()
    soup = BeautifulSoup(html, "lxml")
    slides_container = soup.select_one("div.slides")
    sections = [s for s in slides_container.find_all("section", recursive=False)]

    for i, sect in enumerate(sections):
        slide = prs.slides.add_slide(BLANK)
        # background
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = PAPER
        _layout(slide, sect)
        add_notes(slide, sect)

    prs.save(out_path)
    print("Saved %d slides to %s" % (len(sections), out_path))


if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else OUT_DEFAULT
    main(out)
